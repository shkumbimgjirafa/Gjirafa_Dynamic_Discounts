"""
Builds a visual PowerPoint deck from docs/pricing-algorithms-explained.md.
Plain-language, same content + tone as the doc (post discount-ceiling removal).

Run:  python scripts/build_algorithms_pptx.py
Out:  docs/Pricing Algorithms Explained.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ---------------------------------------------------------------
DARK   = RGBColor(0x1F, 0x2A, 0x37)
NAVY   = RGBColor(0x21, 0x25, 0x29)
BLUE   = RGBColor(0x0D, 0x6E, 0xFD)
GREEN  = RGBColor(0x19, 0x87, 0x54)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
RED    = RGBColor(0xDC, 0x35, 0x45)
LIGHT  = RGBColor(0xF6, 0xF8, 0xFA)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x6C, 0x75, 0x7D)
TEXT   = RGBColor(0x21, 0x25, 0x29)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDE, 0xE2, 0xE6)
FONT   = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, l, t, w, h, fill, line=None, line_w=0.75, rounded=True, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def shape(s, kind, l, t, w, h, fill, line=None):
    shp = s.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, bold, color)."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for (t_, sz, bold, col) in para:
            r = p.add_run(); r.text = t_
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col; r.font.name = FONT
    return tb


def set_text(shp, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (t_, sz, bold, col) in para:
            r = p.add_run(); r.text = t_
            r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col; r.font.name = FONT


def header(s, title, kicker=None):
    rect(s, 0, 0, 13.333, 1.15, DARK, rounded=False)
    rect(s, 0, 1.15, 13.333, 0.07, BLUE, rounded=False)
    text(s, 0.55, 0.12, 12.2, 0.95,
         ([[(kicker, 12, True, RGBColor(0x8A, 0xB4, 0xF8))]] if kicker else []) +
         [[(title, 28, True, WHITE)]],
         anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n):
    text(s, 0.55, 7.02, 8, 0.35, [[("GjirafaMall / Gjirafa50 — Dynamic Pricing", 9, False, MUTED)]])
    text(s, 11.6, 7.02, 1.2, 0.35, [[(str(n), 9, False, MUTED)]], align=PP_ALIGN.RIGHT)


page = [0]
def nextn():
    page[0] += 1
    return page[0]


DIR = {  # direction -> (color, glyph, caption)
    "up":   (GREEN,  "▲", "Price ↑ / protect"),
    "down": (ORANGE, "▼", "Price ↓ / discount"),
    "both": (BLUE,   "⇅", "Either way"),
}


# ============================================================ 1. TITLE
s = slide(DARK)
rect(s, 0, 0, 13.333, 7.5, DARK, rounded=False)
rect(s, 0, 4.95, 13.333, 0.06, BLUE, rounded=False)
text(s, 0.9, 1.9, 11.5, 2.0, [
    [("How the Pricing Tool", 46, True, WHITE)],
    [("Decides a Price", 46, True, WHITE)],
], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.92, 5.15, 11.5, 1.4, [
    [("A plain-language guide to the 10 algorithms,", 20, False, RGBColor(0xCB, 0xD5, 0xE1))],
    [("how their votes combine, and the guardrails that keep prices safe.", 20, False, RGBColor(0xCB, 0xD5, 0xE1))],
], anchor=MSO_ANCHOR.TOP)
text(s, 0.92, 6.7, 11.5, 0.5, [[("Proposals only — humans review & approve; live prices change only via the explicit Push step.",
                                  13, True, RGBColor(0x8A, 0xB4, 0xF8))]])

# ============================================================ 2. BIG PICTURE
s = slide()
header(s, "The big picture", "FROM DATA TO A PROPOSED PRICE")
steps = [
    ("1", "Gather facts", "Price, cost, stock,\nsales over 7–90 days", BLUE),
    ("2", "Advisors vote", "Each algorithm votes\nfor a price or stays silent", RGBColor(0x3D, 0x8B, 0xFD)),
    ("3", "Blend votes", "Weighted average —\nlouder advisors pull more", RGBColor(0x6E, 0xA8, 0xFE)),
    ("4", "Guardrails", "Never below the margin\nfloor; never above shelf", GREEN),
    ("5", "Round price", "Snap to a tidy price\n(if it stays in bounds)", ORANGE),
]
x = 0.45
cw = 2.45
for i, (num, title, desc, col) in enumerate(steps):
    ch = shape(s, MSO_SHAPE.CHEVRON, x, 1.75, cw, 1.5, col)
    try:
        ch.adjustments[0] = 0.26   # shorter arrow point -> wider flat body for the label
    except Exception:
        pass
    # Overlay textbox (own width) so long words like "Guardrails" never break inside the chevron
    text(s, x, 1.92, cw - 0.2, 1.15,
         [[(num + ".", 13, True, RGBColor(0xDC, 0xE7, 0xFF))], [(title, 15, True, WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=1)
    text(s, x + 0.18, 3.45, cw - 0.1, 1.1, [[(line, 12, False, TEXT)] for line in desc.split("\n")],
         align=PP_ALIGN.CENTER)
    x += cw + 0.05
rect(s, 0.45, 5.2, 12.45, 1.4, CARD, line=BORDER)
text(s, 0.75, 5.35, 11.9, 1.15, [
    [("✅  The result is a ", 14, False, TEXT), ("proposal", 14, True, GREEN),
     (", not a live change. A person reviews and approves it.", 14, False, TEXT)],
    [("⚠️  The tool only proposes ", 14, False, TEXT), ("discounts", 14, True, ORANGE),
     (" (at or below the shelf price) — it never raises a price above shelf.", 14, False, TEXT)],
], space=8)
footer(s, nextn())

# ============================================================ 3. GLOSSARY
s = slide()
header(s, "A few words you'll see repeated", "THE VOCABULARY")
terms = [
    ("Shelf price", "The full “list” price. All discounts are measured down from here."),
    ("Current price", "What it actually sells for today (may already be discounted)."),
    ("Cost (PPTCV)", "What the item costs us to buy. Missing cost → the SKU is skipped."),
    ("Margin", "Profit share of the selling price, after cost and after VAT."),
    ("Velocity", "Sales speed — units sold per day, over several windows."),
    ("Days-to-sellout", "At today's pace, how many days until the stock runs out."),
    ("No-sale streak", "How many days in a row the item has sold nothing."),
    ("Discount % (pp)", "How far below shelf price. “pp” = percentage points."),
]
cols, cw, ch, gx, gy = 2, 6.05, 0.93, 0.35, 0.18
x0, y0 = 0.5, 1.45
for i, (term, desc) in enumerate(terms):
    r, c = divmod(i, cols)
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    rect(s, x, y, cw, ch, CARD, line=BORDER)
    rect(s, x, y, 0.12, ch, BLUE, rounded=False)
    text(s, x + 0.28, y + 0.06, cw - 0.35, ch - 0.1, [
        [(term, 14, True, NAVY)],
        [(desc, 11.5, False, RGBColor(0x49, 0x50, 0x57))],
    ], anchor=MSO_ANCHOR.MIDDLE, space=2)
rect(s, 0.5, 6.55, 12.33, 0.6, RGBColor(0xEA, 0xF1, 0xFF), line=RGBColor(0xBF, 0xD6, 0xFF))
text(s, 0.7, 6.6, 12.0, 0.5, [[("Recent sales count more: sales speed weights the last 7 days at 50%, 14 days at 30%, 30 days at 20%.",
                                 12.5, True, RGBColor(0x0B, 0x4F, 0xB8))]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, nextn())

# ============================================================ 4..13 ALGORITHMS
algos = [
    (1, "Sales velocity + inventory forecast", 70, "both",
     "How many days until the stock sells out, at the current selling pace.",
     ["Sells out within ~3 weeks → shave the discount",
      "1.5–3 months of stock → +3pp  ·  3–6 months → +6pp",
      "Over 6 months of stock → +10pp markdown pressure"],
     "No stock, or nothing is selling at all."),
    (2, "New-product protection", 90, "up",
     "How long ago the product launched.",
     ["Launched within the last 90 days → vote full price (0% off)",
      "Protects fresh launches from automatic discounting"],
     "No launch date in the data yet — currently dormant."),
    (3, "Warehouse-stock aging markdown", 50, "down",
     "Stock that has gone quiet but isn't fully dead yet.",
     ["Quiet for ≥1 week → deepen discount +2pp per week",
      "Capped at +12pp on top of today's discount",
      "The gentle nudge for a slowing item"],
     "Still selling this week, or fully dead (that's #7)."),
    (4, "Stockout-risk protection", 80, "up",
     "Days-to-sellout together with how healthy the margin is.",
     ["Sells out within ~14 days AND margin is healthy",
      "→ remove the discount (go to full price)",
      "Why discount what will sell out anyway?"],
     "Not selling fast, or margin sits near the floor."),
    (5, "Price elasticity (heuristic)", 50, "both",
     "Whether deeper discounting actually lifted sales.",
     ["Discounted deeper but sales stayed flat → revert to baseline",
      "Sales jumped clearly → protect the current price"],
     "Recent vs baseline discount is similar; thin data."),
    (6, "Margin-tier prioritization", 40, "both",
     "How much room the margin gives us to play with.",
     ["High margin (≥40%) → can absorb a deeper cut (+3pp)",
      "Thin margin (near the floor) → halve the discount"],
     "Mid-range margins — no opinion."),
    (7, "Dead-stock progressive markdown", 75, "down",
     "Stuck inventory we hold locally: zero sales in 90 days, still in our own warehouse.",
     ["Start at 10% off, add +5pp every two weeks unsold",
      "Deepens toward the margin floor — no ceiling stops it",
      "Only ever marks down, never shrinks a discount"],
     "Any sale in the last 90 days, or stock is supplier-only (none held locally)."),
    (8, "Discount-effectiveness correction", 65, "up",
     "Whether an active discount is actually working.",
     ["≥10% off but 14-day sales flat vs the 90-day baseline",
      "→ halve the discount (stop giving margin away)"],
     "Barely discounted; or truly dead stock (→ #7)."),
    (9, "Velocity-trend momentum", 45, "both",
     "Whether demand is speeding up or slowing down right now.",
     ["Accelerating (≥1.5×) → trim the discount by a third",
      "Slowing (≤0.5×) → add a modest +3pp to hold volume"],
     "Too little history (under 5 units in 90 days)."),
    (10, "Supplier-vs-local stock positioning", 10, "up",
     "Where the stock sits, for locally-stocked fast movers.",
     ["Mostly local stock & selling well → lean to fuller price",
      "Never discounts supplier-only stock that isn't selling"],
     "Not a clearly mostly-local fast mover — a low-weight tie-breaker."),
]
for num, name, weight, direction, watches, does, quiet in algos:
    s = slide()
    header(s, name, f"ALGORITHM {num} OF 10")
    # weight chip (top-right of header)
    chip = rect(s, 10.55, 0.33, 2.25, 0.5, RGBColor(0x2C, 0x3A, 0x4B), line=BLUE, line_w=1)
    set_text(chip, [[("Default weight  ", 11, False, RGBColor(0xCB, 0xD5, 0xE1)),
                     (str(weight), 14, True, WHITE)]], align=PP_ALIGN.CENTER)
    # direction badge (left)
    col, glyph, cap = DIR[direction]
    oval = shape(s, MSO_SHAPE.OVAL, 0.75, 2.05, 2.6, 2.6, col)
    set_text(oval, [[(glyph, 54, True, WHITE)]], align=PP_ALIGN.CENTER)
    text(s, 0.6, 4.8, 2.9, 0.6, [[(cap, 14, True, col)]], align=PP_ALIGN.CENTER)
    # right column blocks
    rx, rw = 4.05, 8.7
    text(s, rx, 1.7, rw, 0.9, [
        [("Watches", 13, True, BLUE)],
        [(watches, 17, False, TEXT)],
    ], space=3)
    rect(s, rx, 2.75, rw, 0.02, BORDER, rounded=False)
    text(s, rx, 2.9, rw, 0.4, [[("What it does", 13, True, BLUE)]])
    text(s, rx + 0.05, 3.3, rw, 2.0,
         [[("•  ", 15, True, col), (b, 15, False, TEXT)] for b in does], space=7)
    qy = 5.75
    rect(s, rx, qy, rw, 0.85, RGBColor(0xF1, 0xF3, 0xF5), line=BORDER)
    text(s, rx + 0.2, qy + 0.05, rw - 0.4, 0.75, [
        [("Stays silent when:  ", 12.5, True, MUTED), (quiet, 12.5, False, RGBColor(0x49, 0x50, 0x57))],
    ], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, nextn())

# ============================================================ 14. VOTES COMBINE
s = slide()
header(s, "How the votes become one price", "CONFLICT RESOLUTION")
rect(s, 0.7, 1.55, 11.9, 1.25, RGBColor(0x10, 0x18, 0x22))
text(s, 0.9, 1.62, 11.5, 1.15, [
    [("vote's pull   =   band weight (0–100)   ×   confidence (0–1)", 19, True, RGBColor(0x8A, 0xB4, 0xF8))],
    [("final price   =   weighted average of all suggested prices, by their pull", 19, True, WHITE)],
], anchor=MSO_ANCHOR.MIDDLE, space=8)
bullets = [
    ("A confident vote from a heavily-weighted algorithm moves the result a lot.", GREEN),
    ("A tentative or low-weighted vote barely nudges it.", BLUE),
    ("Algorithms that stayed silent simply don't participate.", MUTED),
    ("If nobody votes, the price doesn't change.", ORANGE),
]
y = 3.15
for b, col in bullets:
    rect(s, 0.9, y + 0.07, 0.22, 0.22, col)
    text(s, 1.3, y, 11, 0.55, [[(b, 16, False, TEXT)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.72
rect(s, 0.9, 6.35, 11.7, 0.6, RGBColor(0xEA, 0xF1, 0xFF), line=RGBColor(0xBF, 0xD6, 0xFF))
text(s, 1.1, 6.4, 11.3, 0.5, [[("Weights are set per price band — the same advisor can matter a lot for one tier and be switched off for another.",
                                 13, True, RGBColor(0x0B, 0x4F, 0xB8))]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, nextn())

# ============================================================ 15. GUARDRAILS
s = slide()
header(s, "The guardrails (hard limits)", "THREE NON-NEGOTIABLE LIMITS")
# vertical price gauge
gx, gw = 1.2, 2.4
top, bot = 1.7, 6.4
rect(s, gx, top, gw, 0.55, RGBColor(0xAD, 0xB5, 0xBD), rounded=False)          # above shelf (never)
rect(s, gx, top + 0.55, gw, 3.1, RGBColor(0xD1, 0xE7, 0xDD), rounded=False)     # allowed zone
rect(s, gx, top + 3.65, gw, 1.05, RGBColor(0xF8, 0xD7, 0xDA), rounded=False)    # below floor (never)
rect(s, gx, top, gw, bot - top, None, line=NAVY, line_w=1.25)
# labels on the gauge
text(s, gx, top + 0.06, gw, 0.45, [[("✕  above shelf — no raises", 11, True, RGBColor(0x49, 0x50, 0x57))]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, gx, top + 1.6, gw, 0.9, [[("ALLOWED", 15, True, GREEN)], [("discount range", 12, False, GREEN)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, gx, top + 3.95, gw, 0.5, [[("✕  below margin floor", 11, True, RED)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# boundary chips
rect(s, gx + gw - 0.1, top - 0.18, 2.7, 0.45, NAVY); text(s, gx + gw + 0.05, top - 0.16, 2.5, 0.4,
     [[("Shelf-price cap", 12, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, gx + gw - 0.1, top + 3.5, 2.7, 0.45, GREEN); text(s, gx + gw + 0.05, top + 3.52, 2.5, 0.4,
     [[("Margin floor", 12, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
# right-side explanation
rx = 7.0
text(s, rx, 1.7, 5.8, 0.5, [[("Three limits apply after the averaging:", 16, True, NAVY)]])
items = [
    ("1.  Margin floor", "The price can't drop below the level that still earns the band's minimum margin. This is the ONLY limit on how deep a discount can go.", GREEN),
    ("2.  Shelf-price cap", "The price can't go above the full shelf price — the tool proposes discounts, not increases.", NAVY),
    ("3.  Supplier-only dead stock", "Stock only in supplier warehouses with no 90-day sales is never marked below today's price — we don't discount stock we don't hold.", ORANGE),
]
y = 2.3
for t_, d, col in items:
    rect(s, rx, y, 5.8, 0.98, CARD, line=BORDER)
    rect(s, rx, y, 0.12, 0.98, col, rounded=False)
    text(s, rx + 0.28, y + 0.05, 5.4, 0.88, [[(t_, 14, True, col)], [(d, 11.5, False, RGBColor(0x49, 0x50, 0x57))]],
         anchor=MSO_ANCHOR.MIDDLE, space=2)
    y += 1.1
rect(s, rx, y + 0.05, 5.8, 0.95, RGBColor(0xFF, 0xF3, 0xCD), line=RGBColor(0xFF, 0xE0, 0x69))
text(s, rx + 0.2, y + 0.1, 5.4, 0.85, [
    [("No discount ceiling.", 13.5, True, RGBColor(0x8A, 0x6D, 0x00))],
    [("Discounts may go as deep as the margin floor allows — the floor is the sole brake.", 12, False, RGBColor(0x6b, 0x57, 0x00))],
], anchor=MSO_ANCHOR.MIDDLE, space=2)
footer(s, nextn())

# ============================================================ 16. ROUNDING + BANDS
s = slide()
header(s, "Finishing touches", "ROUNDING  &  PRICE BANDS")
# Rounding card
rect(s, 0.5, 1.5, 6.0, 5.2, CARD, line=BORDER)
text(s, 0.8, 1.7, 5.5, 0.5, [[("Psychological rounding", 18, True, NAVY)]])
text(s, 0.8, 2.25, 5.5, 0.6, [[("Snap to a tidy, familiar price — set per band:", 13, False, TEXT)]])
rounding = [
    (".99 / .95 endings", "e.g. 26.99"),
    ("Whole number", "e.g. 27"),
    ("995-style steps", "e.g. 995, 1095"),
    ("…99 whole-currency", "MKD / ALL — e.g. 6199, 9999"),
]
y = 2.95
for a, b in rounding:
    rect(s, 0.8, y, 0.18, 0.18, ORANGE)
    text(s, 1.1, y - 0.07, 5.2, 0.45, [[(a + "  ", 13.5, True, TEXT), (b, 12, False, MUTED)]],
         anchor=MSO_ANCHOR.MIDDLE)
    y += 0.62
rect(s, 0.8, 5.65, 5.4, 0.85, RGBColor(0xF1, 0xF3, 0xF5), line=BORDER)
text(s, 1.0, 5.7, 5.0, 0.75, [[("Rounding is applied only if the rounded price still respects the guardrails — otherwise it's skipped.",
                                 12, False, RGBColor(0x49, 0x50, 0x57))]], anchor=MSO_ANCHOR.MIDDLE)
# Bands card
rect(s, 6.85, 1.5, 6.0, 5.2, CARD, line=BORDER)
text(s, 7.15, 1.7, 5.5, 0.5, [[("Price bands", 18, True, NAVY)]])
text(s, 7.15, 2.25, 5.5, 0.85, [[("Every SKU is sorted into a band by its ", 13, False, TEXT),
                                  ("cost (PPTCV)", 13, True, BLUE),
                                  (" — not its selling price.", 13, False, TEXT)]])
# cost axis with tiers
tiers = ["0–10", "10–50", "50–100", "100–250", "250+"]
tcols = [RGBColor(0x0D,0x6E,0xFD), RGBColor(0x3D,0x8B,0xFD), RGBColor(0x6E,0xA8,0xFE),
         RGBColor(0x9E,0xC5,0xFE), RGBColor(0xCF,0xE2,0xFF)]
tx = 7.15; tw = 1.06
for lab, c in zip(tiers, tcols):
    rect(s, tx, 3.25, tw - 0.06, 0.6, c)
    text(s, tx, 3.27, tw - 0.06, 0.56, [[(lab, 10.5, True, NAVY if c.__str__() in ("9EC5FE","CFE2FF") else WHITE)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx += tw
text(s, 7.15, 3.95, 5.4, 0.4, [[("←  lower cost                              higher cost  →", 11, False, MUTED)]],
     align=PP_ALIGN.CENTER)
text(s, 7.15, 4.55, 5.4, 0.5, [[("Each band carries its own:", 13, True, NAVY)]])
for i, b in enumerate(["Margin floor (the guardrail)", "Rounding style",
                       "On/off + weight for each of the 10 algorithms"]):
    rect(s, 7.35, 5.05 + i*0.5, 0.16, 0.16, GREEN)
    text(s, 7.62, 4.98 + i*0.5, 5.0, 0.45, [[(b, 12.5, False, TEXT)]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, nextn())

# ============================================================ 17. SUMMARY TABLE
s = slide()
header(s, "All 10 algorithms at a glance", "ONE-LINE SUMMARY")
rows = [
    ("1", "Sales velocity + inventory forecast", "down if slow, up if fast", "both"),
    ("2", "New-product protection", "up (full price)", "up"),
    ("3", "Warehouse-stock aging markdown", "down", "down"),
    ("4", "Stockout-risk protection", "up (full price)", "up"),
    ("5", "Price elasticity", "to baseline / hold", "both"),
    ("6", "Margin-tier prioritization", "down if fat, up if thin", "both"),
    ("7", "Dead-stock progressive markdown", "down (local stock only)", "down"),
    ("8", "Discount-effectiveness correction", "up (halve discount)", "up"),
    ("9", "Velocity-trend momentum", "up if accel., down if slowing", "both"),
    ("10", "Supplier-vs-local positioning", "up (toward fuller price)", "up"),
]
tx, ty, tw = 0.5, 1.45, 12.33
rowh = 0.47
# head
rect(s, tx, ty, 0.7, rowh, DARK, rounded=False)
rect(s, tx+0.7, ty, 6.6, rowh, DARK, rounded=False)
rect(s, tx+7.3, ty, 5.03, rowh, DARK, rounded=False)
text(s, tx, ty, 0.7, rowh, [[("#", 12, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, tx+0.85, ty, 6.4, rowh, [[("Algorithm", 12, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, tx+7.45, ty, 4.8, rowh, [[("Pushes price…", 12, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
y = ty + rowh
for i, (num, name, push, d) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF4, 0xF7)
    rect(s, tx, y, tw, rowh, bg, rounded=False)
    col = DIR[d][0]
    text(s, tx, y, 0.7, rowh, [[(num, 12, True, MUTED)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, tx+0.85, y, 6.4, rowh, [[(name, 12.5, False, TEXT)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, tx+7.3, y+0.09, 0.13, rowh-0.18, col, rounded=False)
    text(s, tx+7.55, y, 4.7, rowh, [[(push, 12, True, col)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rowh
rect(s, tx, ty, tw, rowh*(len(rows)+1), None, line=BORDER, line_w=1, rounded=False)
text(s, 0.5, 6.65, 12.3, 0.5, [[("Thresholds shown are tunable defaults. Discounts are measured against the shelf price; margins are computed after VAT.",
                                  11, False, MUTED)]])
footer(s, nextn())

# ============================================================ 18. CLOSING
s = slide(DARK)
rect(s, 0, 0, 13.333, 7.5, DARK, rounded=False)
rect(s, 0, 3.05, 13.333, 0.06, BLUE, rounded=False)
text(s, 0.9, 2.0, 11.5, 1.0, [[("Proposals only.", 40, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.92, 3.3, 11.5, 1.6, [
    [("10 advisors vote  →  weighted blend  →  guardrails  →  rounding.", 20, False, RGBColor(0xCB, 0xD5, 0xE1))],
    [("A human reviews and approves every change. Live prices move only via the explicit Push step.", 18, False, RGBColor(0xCB, 0xD5, 0xE1))],
], anchor=MSO_ANCHOR.TOP, space=10)
text(s, 0.92, 6.7, 11.5, 0.4, [[("Full write-up: docs/pricing-algorithms-explained.md", 13, True, RGBColor(0x8A, 0xB4, 0xF8))]])

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "docs", "Pricing Algorithms Explained.pptx")
prs.save(out)
print("Saved:", out, "|", len(prs.slides._sldIdLst), "slides")
